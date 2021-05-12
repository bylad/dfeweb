import os
import re
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

from django.conf import settings
# os.environ.setdefault('DJANGO_SETTINGS_MODULE', 'dfesite.settings')
# import django
# django.setup()


MEDIA = settings.MEDIA_DIR
MONTHS = ['январь', 'февраль', 'март', 'апрель', 'май', 'июнь', 'июль',
          'август', 'сентябрь', 'октябрь', 'ноябрь', 'декабрь']


def shape_upd(txt_frame, text, fontsize, bold=1):
    txt_frame.clear()
    txt_frame.fit_text()
    p = txt_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(fontsize)
    if bold:
        font.bold = True
    else:
        font.bold = False
    font.color.rgb = RGBColor(55, 96, 146)


def shape_list(shape, current_list, previous_list):
    rup = '↑'
    rdo = '↓'
    # Замена ЗП и стрелок
    if shape.has_text_frame:
        for j in range(9):
            if (shape.name.find(f'nao{j}')) != -1:
                text_frame = shape.text_frame
                if current_list[j] != 0:
                    text_frame.paragraphs[0].runs[0].text = re.sub('\.', ',', str(current_list[j]))
                else:
                    text_frame.paragraphs[0].runs[0].text = "...*"
            if (shape.name.find(f'strelka{j}')) != -1:
                text_frame = shape.text_frame
                text_runs = text_frame.paragraphs[0].runs[0]
                if current_list[j] < previous_list[j]:
                    text_runs.text = rdo
                    text_runs.font.color.rgb = RGBColor(158, 0, 0)
                elif current_list[j] > previous_list[j]:
                    text_runs.text = rup
                    text_runs.font.color.rgb = RGBColor(79, 98, 40)
                elif current_list[j] == previous_list[j]:
                    text_runs.text = ''


def new_pptx(title_date, cur_list, prev_list):
    """ Заполнение презентации актуальными данными
    """
    prs = Presentation(os.path.join(MEDIA, 'salary', 'sample', 'Stat_salary.pptx'))
    date4filename = f"{title_date.year}_01-{str(title_date.month).zfill(2)}"
    regex = re.compile('[яфмаисонд][а-я]+[ьтй]\s*[-]\s*[яфмаисонд][а-я]+[ьтй]\s+\d{4}\s+года?')
    if title_date.month == 1:  # за январь 2020 года
        news_date_text = f"{MONTHS[0]} {title_date.year} года"
    elif title_date.month ==12:  # за 2020 год
        news_date_text = f"{title_date.year} год"
    else:
        news_date_text = f"январь-{MONTHS[title_date.month-1]} {title_date.year} года"

    for slide in prs.slides:
        for sh in slide.shapes:
            # Замена даты в заголовке
            if sh.name == 'zagolovok':
                text_frame = sh.text_frame
                cur_text = sh.text
                cur_text = re.sub(regex, news_date_text, cur_text)
                shape_upd(text_frame, cur_text, 26)
            # Замена года в нижней сноске
            if sh.name == 'snoskayear':
                text_frame = sh.text_frame
                cur_text = sh.text
                snoska_year = int(re.search('\d{4}', cur_text).group())
                if (title_date.year - 1) != snoska_year:
                    c_text = re.sub('\d{4}', str(title_date.year - 1), cur_text)
                    shape_upd(text_frame, c_text, 14, 0)

            shape_list(sh, cur_list, prev_list)

    stat_filename = f'Stat_salary_{date4filename}.pptx'
    path_year = os.path.join(MEDIA, 'salary', f'{title_date.year}')
    if os.path.exists(path_year):
        prs_full_path = os.path.join(path_year, stat_filename)
    else:
        os.mkdir(path_year)
        prs_full_path = os.path.join(path_year, stat_filename)
    print(prs_full_path)
    # if os.path.exists(prs_full_path):
    #     print('Файл с таким именем существует')
    #     return prs_full_path
    prs.save(prs_full_path)
    return prs_full_path

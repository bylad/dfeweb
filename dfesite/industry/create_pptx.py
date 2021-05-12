import os
import re
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

from django.conf import settings
from . import fill_stat_news

MEDIA = settings.MEDIA_DIR
monthe = ['январе', 'феврале', 'марте', 'апреле', 'мае', 'июне', 'июле',
          'августе', 'сентябре', 'октябре', 'ноябре', 'декабре']
months = ['январь', 'февраль', 'март', 'апрель', 'май', 'июнь', 'июль',
          'август', 'сентябрь', 'октябрь', 'ноябрь', 'декабрь']

def date_string(txt):
    """ Поиск по регулярному выражению.
        Например, из строки "О производстве в январе - мае 2019 года"
        будет выделена подстрока: "январе - мае 2019"
    """
    dt = re.compile("[яфмаисонд]([а-я]+[е])( *)(.)( *)[яфмаисонд]([а-я]+[е])( |.)\d{4}")
    match = re.search(dt, txt)
    if match:
        return match.group()


def shape_upd(txt_frame, text, fontsize):
    txt_frame.clear()
    txt_frame.fit_text()
    p = txt_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(fontsize)
    font.bold = True
    font.color.rgb = RGBColor(55, 96, 146)


def month_year(txt, d):
    if len(d) == 1:  # в заголовке указан только год (статистика за весь год)
        new_date = f'{str(d[0])}'
        new_month = f'{str(d[0])} в %\n{str(d[0]-1)} г.'
        fname_date = f'{str(d[0])}-01-12'
        return new_date, new_month, fname_date
    elif len(d) == 2:  # в заголовке указан год и январь
        new_date = f'{monthe[d[1]-1]} {str(d[0])}'
        new_month = f'{months[d[1]-1]} {str(d[0])} в %\n{months[d[1]-1]} {str(d[0]-1)} г.'
        fname_date = f'{str(d[0])}-01-01'
        return new_date, new_month, fname_date
    # в заголовке указан год и период с января по указанный месяц
    new_date = f'{monthe[0]}-{monthe[d[2]-1]} {str(d[0])}'
    new_month = f'{months[0][:3]}-{months[d[2]-1]} {str(d[0])} в %\n{months[0][:3]}-{months[d[2]-1]} {str(d[0]-1)} г.'
    fname_date = f'{str(d[0])}-01-{"{:02d}".format(d[2])}'

    print('new_date, new_month, fname_date')
    print('-------------------------------')
    print(new_date, new_month, fname_date)
    return new_date, new_month, fname_date


def new_pptx(production_list, idx, news_title):
    """ Заполнение презентации актуальными данными
        - zagolovok: замена указанных в заголовке месяца (года)
        - mes: замена указанных в надписи месяца (года)
        - index: замена индекса в %
        - prom0-7: замена соответствующих показателей
    """
    location = os.path.join(MEDIA, 'industry', 'sample', 'Stat_industry.pptx')
    prs = Presentation(location)
    date_list = fill_stat_news.date_int(news_title)
    newdate, newmonth, date4filename = month_year(news_title, date_list)
    for slide in prs.slides:
        for shape in slide.shapes:
            # Замена даты в заголовке
            if shape.name == 'zagolovok':
                text_frame = shape.text_frame
                cur_text = shape.text
                search_str = date_string(shape.text)
                if len(newdate) < 5:  # только год, => меняем года на году
                    search_str = search_str + ' года'
                    newdate = newdate + ' году'
                new_text = cur_text.replace(search_str, newdate)
                shape_upd(text_frame, new_text, 26)
            # Замена даты в индексе
            if shape.name == 'mes':
                text_frame = shape.text_frame
                shape_upd(text_frame, newmonth, 16)
            if shape.name == 'index':
                text_frame = shape.text_frame
                indx = re.sub('\.', ',', str(idx)) + '%'
                shape_upd(text_frame, indx, 40)
            for k in range(8):
                if shape.name == 'prom' + str(k):
                    text_frame = shape.text_frame
                    if production_list[k] == '0':
                        shape_upd(text_frame, '-', 40)
                    else:
                        shape_upd(text_frame, re.sub('\.', ',', production_list[k]), 40)

    stat_filename = f'Stat_industry_{date4filename}.pptx'
    path_year = os.path.join(MEDIA, 'industry', f'{date_list[0]}')
    if os.path.exists(path_year):
        prs_full_path = os.path.join(path_year, stat_filename)
    else:
        os.mkdir(path_year)
        prs_full_path = os.path.join(path_year, stat_filename)

    print(prs_full_path)
    # if os.path.exists(prs_full_path):
    #     print('Файл с таким именем существует')
    #     return prs_full_path
    prs.save(prs_full_path)
    return prs_full_path

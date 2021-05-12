import os
import re
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import dateparser

from django.conf import settings

MEDIA = settings.MEDIA_DIR
MONTHS = ['января', 'февраля', 'марта', 'апреля', 'мая', 'июня', 'июля',
          'августа', 'сентября', 'октября', 'ноября', 'декабря']

def cut_date(txt):
    """
    Из заголовка вычленяем дату. Например из текста (txt):
        "Средние  цены на отдельные потребительские товары (услуги)
         по Ненецкому автономному округу на 13 января 2020 года"
    получим "13 января 2020". Ф-я возвращает эту строку в формате дата.
    """
    regex = re.compile('\d{1,2}\s[яфмаисонд][а-я]+[а|я]\s\d{4}')
    return dateparser.parse(re.search(regex, txt).group())


def shape_upd(txt_frame, newtext, fontsize, align='LEFT'):
    txt_frame.clear()
    txt_frame.fit_text()
    p = txt_frame.paragraphs[0]
    run = p.add_run()
    run.text = newtext
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(fontsize)
    font.bold = True
    font.color.rgb = RGBColor(55, 96, 146)
    if align == 'RIGHT':  # RIGHT => форматируем значения цен
        p.alignment = PP_ALIGN.RIGHT
        font.color.rgb = RGBColor(37, 64, 97)


def pptx_table(table, current_list, previous_list, m, n):
    # table0=14, table1=3, table2=8
    # 1-ая строка в таблицах PPTX содержит имя столбца => отсчет не с 0, а с 1
    row = 1
    for i in range(m, n):
        if current_list[i] == previous_list[i]:
            # table.cell(row, 1).text = str(current_list[i]).zfill(2)
            table.cell(row, 2).text = ''
        elif current_list[i] > previous_list[i]:
            # table.cell(row, 1).text = str(current_list[i]).zfill(2)
            table.cell(row, 2).text = '↑'
            table.cell(row, 2).text_frame.paragraphs[0].font.color.rgb = RGBColor(158, 0, 0)  # Red
        elif current_list[i] < previous_list[i]:
            # table.cell(row, 1).text = str(current_list[i]).zfill(2)
            table.cell(row, 2).text = '↓'
            table.cell(row, 2).text_frame.paragraphs[0].font.color.rgb = RGBColor(79, 98, 40)  # Green
        # округляем вещ.число до 2х знаков после запятой, затем форматируем (добавляет 0)
        shape_upd(table.cell(row, 1).text_frame, format(round(current_list[i], 2), '.2f'), 16, 'RIGHT')
        table.cell(row, 2).text_frame.paragraphs[0].font.size = Pt(16)
        table.cell(row, 2).text_frame.paragraphs[0].font.bold = True

        row += 1


def new_pptx(news_title, current_price_list, previous_price_list, previous_pub_date):
    """ Заполнение презентации актуальными данными
        - zagolovok: замена даты в заголовке
        - prom0-7: замена соответствующих показателей
    """
    prs = Presentation(os.path.join(MEDIA, 'price', 'sample', 'Stat_price.pptx'))
    regex = re.compile('\d{1,2}\s[яфмаисонд][а-я]+[а|я]\s\d{4}')
    news_date_text = re.search(regex, news_title).group()
    news_date = dateparser.parse(news_date_text)
    date4filename = f"{news_date.year}-{str(news_date.month).zfill(2)}-{str(news_date.day).zfill(2)}"
    for slide in prs.slides:
        for shape in slide.shapes:
            # Замена даты в заголовке
            if shape.name == 'zagolovok':
                text_frame = shape.text_frame
                cur_text = shape.text
                cur_text = re.sub(regex, news_date_text, cur_text)
                shape_upd(text_frame, cur_text, 26)

            if shape.name == 'snoskadate':
                text_frame = shape.text_frame
                shape_upd(text_frame, previous_pub_date, 10, 0)

            if shape.name == 'table0':
                pptx_table(shape.table, current_price_list, previous_price_list, 0, 14)
            if shape.name == 'table1':
                pptx_table(shape.table, current_price_list, previous_price_list, 14, 17)
            if shape.name == 'table2':
                pptx_table(shape.table, current_price_list, previous_price_list, 17, 25)

    stat_filename = f'Stat_price_{date4filename}.pptx'
    path_year = os.path.join(MEDIA, 'price', f'{news_date.year}')
    
    if os.path.exists(path_year):
        prs_full_path = os.path.join(path_year, stat_filename)
    else:
        os.mkdir(path_year)
        prs_full_path = os.path.join(path_year, stat_filename)
    print(prs_full_path)
    # if os.path.exists(prs_full_path):
    #     print('Файл с таким именем существует')
    #     return prs_full_path
    prs.save(prs_full_path)
    return prs_full_path

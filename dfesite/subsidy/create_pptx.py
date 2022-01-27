import re
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN

from django.conf import settings

MEDIA = settings.MEDIA_DIR


def shape_upd(txt_frame, newtext, fontsize, align='LEFT'):
    txt_frame.clear()
    txt_frame.fit_text()
    p = txt_frame.paragraphs[0]
    run = p.add_run()
    run.text = newtext
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(fontsize)
    font.bold = True
    font.color.rgb = RGBColor(55, 96, 146)
    if align == 'RIGHT':  # RIGHT => форматируем значения цен
        p.alignment = PP_ALIGN.RIGHT
        font.color.rgb = RGBColor(37, 64, 97)


def new_pptx(period, data_dict):
    """ Заполнение презентации актуальными данными
        period_txt: строка вида '9 месяцев 2021 года'
        data_dict: словарь с данными
    """
    prs = Presentation(Path(MEDIA, 'subsidy', 'sample', 'Stat_subsidy.pptx'))
    regex = re.compile(r"\d\s*(месяца|месяцев)\s*\d{4}\s*года")  # '9 месяцев 2021 года'

    for slide in prs.slides:
        for shape in slide.shapes:
            # Замена даты в заголовке
            if shape.name == 'zagolovok':
                text_frame = shape.text_frame
                cur_text = shape.text
                cur_text = re.sub(regex, period[0], cur_text)
                shape_upd(text_frame, cur_text, 26)

            # Замена данных
            for key, value in data_dict.items():
                if shape.name == key:
                    text_frame = shape.text_frame
                    if key == 'subsidy_percent':
                        value_text = f"{str(value).replace('.', ',')}%"
                    elif key in ('subsidy_accrued', 'benefit_funds', 'benefit_compensate'):
                        value_text = f"{str(value).replace('.', ',')} млн"
                    else:
                        value_text = f"{str(value).replace('.', ',')}"
                    shape_upd(text_frame, value_text, 26)

    stat_filename = f'Stat_subsidy_{period[1]}.pptx'
    year = re.search(r"\d{4}", period[0]).group()
    path_year = Path(MEDIA, 'subsidy', f'{year}')
    if path_year.exists():
        prs_full_path = Path(path_year, stat_filename)
    else:
        path_year.parent.mkdir(exist_ok=True, parent=True)
        path_year.touch()
        prs_full_path = Path(path_year, stat_filename)
    prs.save(prs_full_path)
    return str(prs_full_path)

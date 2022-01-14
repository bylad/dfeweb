import os
import re
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN

from django.conf import settings

MEDIA = settings.MEDIA_DIR
MONTH = ['январь', 'февраль', 'март', 'апрель', 'май', 'июнь', 'июль',
         'август', 'сентябрь', 'октябрь', 'ноябрь', 'декабрь']


def shape_upd(txt_frame, newtext, fontsize, align='LEFT'):
    txt_frame.clear()
    txt_frame.fit_text()
    p = txt_frame.paragraphs[0]
    run = p.add_run()
    run.text = newtext
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(fontsize)
    font.bold = True
    font.color.rgb = RGBColor(55, 96, 146)
    if align == 'RIGHT':  # RIGHT => форматируем значения цен
        p.alignment = PP_ALIGN.RIGHT
        font.color.rgb = RGBColor(37, 64, 97)


def data_handle(migr, zags):
    """ Возвращает 2 списка
            migr_list: 0 - arrived, 1 - departed, 2 - gain
            zags_list: 0 - born, 1 - died, 2 - wedd, 3 - divorce, 4 - gain
    """
    if migr['gain'] > 0:
        m_list = [str(migr['arrived']), str(migr['departed']), f"+{str(migr['gain'])}"]
    else:
        m_list = [str(migr['arrived']), str(migr['departed']), str(migr['gain'])]

    zags_compare = zags['born'] - zags['died']
    if zags_compare > 0:
        zags_gain = f"+{str(zags_compare)}"
    else:
        zags_gain = f"{str(zags_compare)}"

    z_list = [str(zags['born']), str(zags['died']), str(zags['wedd']), str(zags['divorce']), zags_gain]

    return m_list, z_list


def new_pptx(head_date, migr_dict, zags_dict):
    """ Заполнение презентации актуальными данными
        head_date: 0 - int_year, 1 - int_month, 2 - str_monthe
        zagolovok: замена даты в заголовке
        migr, zags: замена соответствующих показателей
    """
    prs = Presentation(Path(MEDIA, 'population', 'sample', 'Stat_population.pptx'))
    regex = re.compile(r'\s[яфмаисонд][а-я]+[ьтй]\s\d{4}')
    new_date = f" {MONTH[head_date[1]]} {head_date[0]}"
    date4filename = f"{str(head_date[0])}-{str(head_date[1] + 1).zfill(2)}"
    migr_list, zags_list = data_handle(migr_dict, zags_dict)
    for slide in prs.slides:
        for shape in slide.shapes:
            # Замена даты в заголовке
            if shape.name == 'zagolovok':
                text_frame = shape.text_frame
                cur_text = shape.text
                cur_text = re.sub(regex, new_date, cur_text)
                shape_upd(text_frame, cur_text, 26)

            # Замена данных
            if shape.name == 'migr_arrived':
                text_frame = shape.text_frame
                shape_upd(text_frame, migr_list[0], 40)
            if shape.name == 'migr_departed':
                text_frame = shape.text_frame
                shape_upd(text_frame, migr_list[1], 40)
            if shape.name == 'migr_gain':
                text_frame = shape.text_frame
                shape_upd(text_frame, migr_list[2], 40)

            if shape.name == 'zags_born':
                text_frame = shape.text_frame
                shape_upd(text_frame, zags_list[0], 40)
            if shape.name == 'zags_died':
                text_frame = shape.text_frame
                shape_upd(text_frame, zags_list[1], 40)
            if shape.name == 'zags_wedd':
                text_frame = shape.text_frame
                shape_upd(text_frame, zags_list[2], 40)
            if shape.name == 'zags_divorce':
                text_frame = shape.text_frame
                shape_upd(text_frame, zags_list[3], 40)
            if shape.name == 'zags_gain':
                text_frame = shape.text_frame
                shape_upd(text_frame, zags_list[4], 40)

    stat_filename = f'Stat_population_{date4filename}.pptx'
    path_year = Path(MEDIA, 'population', f'{str(head_date[0])}')
    if os.path.exists(path_year):
        prs_full_path = Path(path_year, stat_filename)
    else:
        os.mkdir(path_year)
        prs_full_path = Path(path_year, stat_filename)
    prs.save(prs_full_path)
    return str(prs_full_path)
